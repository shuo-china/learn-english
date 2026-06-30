from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
IMG = ROOT / "images"
OUT = ROOT / "forest_vocabulary_adventure.pptx"

W, H = Inches(13.333333), Inches(7.5)
GREEN = RGBColor(70, 99, 48)
DARK = RGBColor(55, 72, 42)
GOLD = RGBColor(166, 137, 67)
CREAM = RGBColor(255, 248, 224)
BLUE = RGBColor(82, 145, 158)
ORANGE = RGBColor(196, 125, 49)


def set_font(run, size, bold=False, color=DARK, name="Microsoft YaHei"):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(slide, text, x, y, w, h, size=22, bold=False, color=DARK, align=PP_ALIGN.CENTER):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.clear()
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    set_font(r, size, bold, color)
    return box


def add_multiline(slide, lines, x, y, w, h, size=15, color=DARK, bold_first=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = line
        set_font(r, size, bold_first and i == 0, color)
    return box


def add_round_rect(slide, x, y, w, h, fill=CREAM, line=RGBColor(213, 185, 115), transparency=18):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.fill.transparency = transparency
    shp.line.color.rgb = line
    shp.line.width = Pt(1.2)
    return shp


def add_icon_circle(slide, text, x, y, color, size=0.46):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.color.rgb = RGBColor(255, 255, 255)
    shp.line.width = Pt(1)
    add_text(slide, text, x, y + 0.01, size, size - 0.02, size=18, bold=True, color=RGBColor(255, 255, 255))


def title(slide, text, y=0.34, size=27):
    add_text(slide, f"❧ {text} ❧", 2.0, y, 9.4, 0.45, size=size, bold=True, color=GREEN)


def background(slide, n):
    slide.shapes.add_picture(str(IMG / f"slide{n:02d}_bg.png"), 0, 0, width=W, height=H)


prs = Presentation()
prs.slide_width = W
prs.slide_height = H
blank = prs.slide_layouts[6]


# Slide 1
s = prs.slides.add_slide(blank)
background(s, 1)
add_text(s, "把单词练习", 2.05, 2.12, 4.6, 0.6, size=31, bold=True, color=GREEN)
add_text(s, "变成一场森林探险", 1.55, 2.78, 5.65, 0.7, size=30, bold=True, color=GREEN)
add_round_rect(s, 2.45, 3.82, 3.35, 0.36, fill=RGBColor(99, 126, 69), line=RGBColor(99, 126, 69), transparency=0)
add_text(s, "Vocabulary | 儿童英语词汇背诵与拼写练习工具", 2.55, 3.85, 3.15, 0.23, size=9, bold=False, color=RGBColor(255, 249, 225))


# Slide 2
s = prs.slides.add_slide(blank)
background(s, 2)
title(s, "孩子练单词，常卡在三件事", y=0.38, size=24)
for text, x in [("记得慢", 2.25), ("写不准", 5.15), ("坚持难", 8.02)]:
    add_text(s, text, x, 4.6, 1.3, 0.35, size=17, bold=True, color=DARK)
add_round_rect(s, 4.95, 5.75, 3.45, 0.56, fill=CREAM, transparency=8)
add_text(s, "需要一个更轻、更有反馈的练习场景", 5.08, 5.88, 3.2, 0.24, size=13, bold=True, color=DARK)


# Slide 3
s = prs.slides.add_slide(blank)
background(s, 3)
title(s, "双模式学习路径", y=0.36, size=24)
rows = [
    ("▣", "背诵：点击揭示中文释义，\n听发音加深记忆"),
    ("✎", "拼写：根据中文提示输入英文，\n逐段校验"),
    ("↝", "顺序 / 打乱：兼顾课本\n复习与随机巩固"),
]
for i, (ico, body) in enumerate(rows):
    y = 1.35 + i * 1.05
    add_text(s, ico, 1.28, y + 0.16, 0.45, 0.38, size=20, bold=True, color=GREEN)
    add_multiline(s, [body], 2.0, y + 0.06, 2.2, 0.62, size=12, color=DARK)


# Slide 4
s = prs.slides.add_slide(blank)
background(s, 4)
title(s, "即时反馈，让错误变成线索", y=0.38, size=23)
cards = [
    ("✓", GREEN, "答对：锁定正确片段\n并播放正确音效"),
    ("×", RGBColor(207, 108, 62), "答错：标红提醒，连续错误\n自动显示答案"),
    ("♪", BLUE, "发音：随时重听，建立\n声音与拼写连接"),
]
for i, (ico, c, body) in enumerate(cards):
    y = 1.42 + i * 1.06
    add_icon_circle(s, ico, 0.8, y + 0.09, c)
    add_multiline(s, [body], 1.62, y, 2.0, 0.72, size=12, color=DARK, bold_first=False)


# Slide 5
s = prs.slides.add_slide(blank)
background(s, 5)
title(s, "词书驱动，内容可持续扩展", y=0.38, size=23)
features = [
    ("≡", GREEN, "支持多个词书"),
    ("★", RGBColor(218, 169, 57), "重点词可标记"),
    ("▰", BLUE, "适合故事词汇、课文词表、\n家庭自定义练习"),
]
for i, (ico, c, body) in enumerate(features):
    y = 1.58 + i * 1.04
    add_icon_circle(s, ico, 0.78, y + 0.02, c)
    add_multiline(s, [body], 1.55, y, 2.7, 0.52, size=13, color=DARK)


# Slide 6
s = prs.slides.add_slide(blank)
background(s, 6)
title(s, "产品价值", y=0.42, size=24)
values = [
    ("♥", GREEN, "给孩子：", "更轻松的练习体验"),
    ("●", ORANGE, "给家长：", "更清晰的学习反馈"),
    ("◆", BLUE, "给老师：", "更灵活词表工具"),
]
for i, (ico, c, head, body) in enumerate(values):
    y = 1.58 + i * 0.72
    add_icon_circle(s, ico, 6.25, y + 0.02, c, size=0.38)
    add_text(s, head, 6.82, y, 1.0, 0.35, size=15, bold=True, color=c, align=PP_ALIGN.LEFT)
    add_text(s, body, 7.72, y, 2.35, 0.35, size=12, bold=False, color=DARK, align=PP_ALIGN.LEFT)
add_text(s, "让每一次拼写，都成为向前一步", 5.9, 5.95, 4.15, 0.35, size=16, bold=True, color=DARK)


prs.save(OUT)
print(OUT)
